"""PPT 文档解析器"""
from pptx import Presentation
from typing import List, Dict, Any
from pathlib import Path


class PPTParser:
    """PPTX 文档解析器"""
    
    def parse(self, file_path: str) -> Dict[str, Any]:
        """解析 PPTX 文件，返回结构化数据
        
        Args:
            file_path: PPTX 文件路径
            
        Returns:
            包含幻灯片信息的字典
        """
        prs = Presentation(file_path)
        
        slides_data = []
        for idx, slide in enumerate(prs.slides, 1):
            slide_data = {
                "slide_number": idx,
                "title": self._extract_title(slide),
                "text_content": self._extract_text(slide),
                "images": self._extract_images(slide, idx),
                "structure": self._extract_structure(slide),
            }
            slides_data.append(slide_data)
        
        return {
            "filename": Path(file_path).name,
            "total_slides": len(slides_data),
            "slides": slides_data
        }
    
    def extract_text(self, file_path: str, file_id: str = None) -> str:
        """提取 PPT 的纯文本内容
        
        Args:
            file_path: PPTX 文件路径
            file_id: 文档ID（可选，用于嵌入元数据标记）
            
        Returns:
            所有幻灯片文本内容（用换行符分隔，每张幻灯片前添加 [FILE:{file_id}][SLIDE:N] 标记）
        """
        prs = Presentation(file_path)
        texts = []
        
        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_texts = []
            
            # 构建元数据标记
            if file_id:
                slide_texts.append(f"[FILE:{file_id}][SLIDE:{slide_idx}]")
            else:
                slide_texts.append(f"[SLIDE:{slide_idx}]")
            
            # 提取标题
            title = self._extract_title(slide)
            if title:
                slide_texts.append(title)
            # 提取内容
            content = self._extract_text(slide)
            if content:
                slide_texts.append(content)
            
            if slide_texts:
                texts.append("\n".join(slide_texts))
        
        return "\n\n".join(texts)
    
    def extract_pages(self, file_path: str, file_id: str = None) -> List[Dict[str, Any]]:
        """提取 PPT 的幻灯片内容（页级三元库）
        
        Args:
            file_path: PPTX 文件路径
            file_id: 文档ID
            
        Returns:
            幻灯片列表，每个元素包含 page_index (slide_number) 和 content
        """
        prs = Presentation(file_path)
        pages = []
        
        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_text = ""
            
            # 提取标题
            title = self._extract_title(slide)
            if title:
                slide_text += title + "\n"
            
            # 提取内容
            content = self._extract_text(slide)
            if content:
                slide_text += content
            
            if slide_text.strip():
                pages.append({
                    "page_index": slide_idx,
                    "content": slide_text.strip()
                })
        
        return pages
    
    def _extract_title(self, slide) -> str:
        """提取幻灯片标题"""
        # 尝试从占位符获取标题
        for shape in slide.shapes:
            if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                if shape.placeholder_format.idx == 0:  # 标题占位符通常是 idx=0
                    if hasattr(shape, "text") and shape.text:
                        return shape.text.strip()
        
        # 如果没有找到占位符标题，尝试查找第一个文本框（通常是标题）
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text = shape.text.strip()
                if text and len(text) < 200:  # 标题通常较短
                    return text
        
        return ""
    
    def _extract_text(self, slide) -> str:
        """提取幻灯片文本内容"""
        texts = []
        title_found = False
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text = shape.text.strip()
                if text:
                    # 跳过标题（已在标题中提取）
                    if not title_found and len(text) < 200:
                        title_found = True
                        # 如果已经在标题中，跳过
                        continue
                    texts.append(text)
        
        return "\n".join(texts)
    
    def _extract_images(self, slide, slide_number: int) -> List[Dict[str, Any]]:
        """提取幻灯片中的图片元信息
        
        Args:
            slide: 幻灯片对象
            slide_number: 幻灯片编号
            
        Returns:
            图片信息列表
        """
        images = []
        for idx, shape in enumerate(slide.shapes):
            if hasattr(shape, "image"):
                image_info = {
                    "image_id": f"slide_{slide_number}_img_{idx}",
                    "position": {
                        "left": int(shape.left),
                        "top": int(shape.top),
                        "width": int(shape.width),
                        "height": int(shape.height),
                    },
                    "alt_text": getattr(shape, "name", f"Image {idx + 1}")
                }
                images.append(image_info)
        return images
    
    def _extract_structure(self, slide) -> Dict[str, Any]:
        """提取幻灯片结构信息"""
        return {
            "layout": getattr(slide.slide_layout, "name", "Unknown"),
            "shapes_count": len(slide.shapes)
        }
    
